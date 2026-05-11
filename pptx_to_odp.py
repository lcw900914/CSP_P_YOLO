"""
PPTX → ODP 文字轉換器（靜態文字 + 表格，圖片略過）
ODF 規範：所有樣式必須先在 automatic-styles 宣告，再以 style-name 引用
"""
import zipfile, re
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

PPTX   = r'C:/Users/lcw/Desktop/AI_Code/CSP_YOLO/CSP_YOLO(二版).pptx'
OUTPUT = r'C:/Users/lcw/Desktop/AI_Code/CSP_YOLO/CSP_YOLO(二版).odp'

prs = Presentation(PPTX)
SW  = prs.slide_width  / 360000.0
SH  = prs.slide_height / 360000.0

# ── 樣式登記表（統一收集，避免重複） ──────────────────────────────
para_styles = {}   # key=(align,size_pt) -> style_name "P1","P2"...
text_styles = {}   # key=(size_pt,bold,italic,color_hex) -> "T1","T2"...

def para_style_name(align, size):
    key = (align, round(size, 1))
    if key not in para_styles:
        para_styles[key] = f'P{len(para_styles)+1}'
    return para_styles[key]

def text_style_name(size, bold, italic, color):
    key = (round(size, 1), bool(bold), bool(italic), color or '')
    if key not in text_styles:
        text_styles[key] = f'T{len(text_styles)+1}'
    return text_styles[key]

# ── 工具 ────────────────────────────────────────────────────────────
def esc(s):
    return (str(s).replace('&','&amp;').replace('<','&lt;')
                  .replace('>','&gt;').replace('"','&quot;'))

def rgb_str(color):
    try:
        return f'#{color.rgb}'
    except Exception:
        return None

def run_size(run):
    try:
        if run.font.size:
            return run.font.size.pt
    except Exception:
        pass
    return None

def para_size(para):
    for r in para.runs:
        s = run_size(r)
        if s:
            return s
    return 18.0

def fo_align(para):
    try:
        a = para.alignment
        if a == PP_ALIGN.CENTER:   return 'center'
        if a == PP_ALIGN.RIGHT:    return 'end'
        if a == PP_ALIGN.JUSTIFY:  return 'justify'
    except Exception:
        pass
    return 'start'

# ── 第一遍：收集所有樣式 ─────────────────────────────────────────
def collect_styles_para(para, default_size=18.0):
    size  = para_size(para) or default_size
    align = fo_align(para)
    para_style_name(align, size)
    for run in para.runs:
        s     = run_size(run) or size
        bold  = bool(run.font.bold)
        ital  = bool(run.font.italic)
        color = rgb_str(run.font.color)
        text_style_name(s, bold, ital, color)

def collect_styles():
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for para in cell.text_frame.paragraphs:
                                collect_styles_para(para)
                elif shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        collect_styles_para(para)
            except Exception:
                pass

collect_styles()

# ── 段落 XML ────────────────────────────────────────────────────────
def para_xml(para, default_size=18.0):
    size  = para_size(para) or default_size
    align = fo_align(para)
    psn   = para_style_name(align, size)

    runs_xml = ''
    for run in para.runs:
        txt = esc(run.text)
        if not txt:
            continue
        s     = run_size(run) or size
        bold  = bool(run.font.bold)
        ital  = bool(run.font.italic)
        color = rgb_str(run.font.color)
        tsn   = text_style_name(s, bold, ital, color)
        runs_xml += f'<text:span text:style-name="{tsn}">{txt}</text:span>'

    if not runs_xml:
        return f'<text:p text:style-name="{psn}"> </text:p>'
    return f'<text:p text:style-name="{psn}">{runs_xml}</text:p>'

# ── 文字框 XML ──────────────────────────────────────────────────────
def textbox_xml(shape):
    x = shape.left   / 360000.0
    y = shape.top    / 360000.0
    w = shape.width  / 360000.0
    h = shape.height / 360000.0
    body = ''.join(para_xml(p) for p in shape.text_frame.paragraphs)
    return (
        f'<draw:frame draw:layer="layout" '
        f'svg:x="{x:.4f}cm" svg:y="{y:.4f}cm" '
        f'svg:width="{w:.4f}cm" svg:height="{h:.4f}cm">'
        f'<draw:text-box>{body}</draw:text-box>'
        f'</draw:frame>'
    )

# ── 表格 XML ────────────────────────────────────────────────────────
def table_xml(shape):
    tbl = shape.table
    x   = shape.left   / 360000.0
    y   = shape.top    / 360000.0
    w   = shape.width  / 360000.0
    h   = shape.height / 360000.0
    nc  = len(tbl.columns)
    col_w = w / nc if nc else w

    cols_xml = ''.join(
        f'<table:table-column table:style-name="col{i+1}"/>'
        for i in range(nc)
    )
    rows_xml = ''
    for row in tbl.rows:
        cells_xml = ''
        for cell in row.cells:
            cell_content = ''.join(
                para_xml(p) for p in cell.text_frame.paragraphs
            )
            cells_xml += (
                f'<table:table-cell table:style-name="ce1" '
                f'office:value-type="string">'
                f'{cell_content}'
                f'</table:table-cell>'
            )
        rows_xml += f'<table:table-row>{cells_xml}</table:table-row>'

    return (
        f'<draw:frame draw:layer="layout" '
        f'svg:x="{x:.4f}cm" svg:y="{y:.4f}cm" '
        f'svg:width="{w:.4f}cm" svg:height="{h:.4f}cm">'
        f'<table:table>{cols_xml}{rows_xml}</table:table>'
        f'</draw:frame>'
    )

# ── 投影片 XML ──────────────────────────────────────────────────────
def slide_xml(idx, slide):
    shapes_xml = ''
    for shape in slide.shapes:
        try:
            if shape.has_table:
                shapes_xml += table_xml(shape)
            elif shape.has_text_frame:
                shapes_xml += textbox_xml(shape)
        except Exception as e:
            print(f'  [WARN] slide {idx+1} shape skip: {e}')
    return (
        f'<draw:page draw:name="page{idx+1}" draw:style-name="dp1" '
        f'draw:master-page-name="Default">'
        f'{shapes_xml}'
        f'</draw:page>'
    )

# ── 自動樣式 XML ────────────────────────────────────────────────────
NS = (
    'xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" '
    'xmlns:draw="urn:oasis:names:tc:opendocument:xmlns:drawing:1.0" '
    'xmlns:presentation="urn:oasis:names:tc:opendocument:xmlns:presentation:1.0" '
    'xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0" '
    'xmlns:table="urn:oasis:names:tc:opendocument:xmlns:table:1.0" '
    'xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0" '
    'xmlns:fo="urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0" '
    'xmlns:svg="urn:oasis:names:tc:opendocument:xmlns:svg-compatible:1.0" '
    'xmlns:xlink="http://www.w3.org/1999/xlink" '
    'xmlns:dc="http://purl.org/dc/elements/1.1/" '
    'xmlns:meta="urn:oasis:names:tc:opendocument:xmlns:meta:1.0" '
    'xmlns:number="urn:oasis:names:tc:opendocument:xmlns:datastyle:1.0" '
    'office:version="1.2"'
)

def build_auto_styles(nc_max=20):
    out = []
    # 投影片頁面樣式
    out.append(
        '<style:style style:name="dp1" style:family="drawing-page">'
        '<style:drawing-page-properties draw:fill="solid" draw:fill-color="#ffffff"/>'
        '</style:style>'
    )
    # 欄寬（最多 nc_max 欄）
    for i in range(nc_max):
        out.append(
            f'<style:style style:name="col{i+1}" style:family="table-column">'
            f'<style:table-column-properties style:column-width="3cm"/>'
            f'</style:style>'
        )
    # 儲存格樣式
    out.append(
        '<style:style style:name="ce1" style:family="table-cell">'
        '<style:table-cell-properties fo:border="0.05pt solid #aaaaaa" fo:padding="0.1cm"/>'
        '</style:style>'
    )
    # 段落樣式
    for (align, size), name in para_styles.items():
        out.append(
            f'<style:style style:name="{name}" style:family="paragraph">'
            f'<style:paragraph-properties fo:text-align="{align}"/>'
            f'<style:text-properties fo:font-size="{size:.1f}pt" '
            f'fo:font-family="Microsoft JhengHei"/>'
            f'</style:style>'
        )
    # 文字樣式
    for (size, bold, italic, color), name in text_styles.items():
        props = [f'fo:font-size="{size:.1f}pt"',
                 'fo:font-family="Microsoft JhengHei"']
        if bold:
            props.append('fo:font-weight="bold"')
        if italic:
            props.append('fo:font-style="italic"')
        if color:
            props.append(f'fo:color="{color}"')
        out.append(
            f'<style:style style:name="{name}" style:family="text">'
            f'<style:text-properties {" ".join(props)}/>'
            f'</style:style>'
        )
    return '\n'.join(out)

def build_content(slides_xml):
    auto = build_auto_styles()
    return f"""<?xml version="1.0" encoding="UTF-8"?>
<office:document-content {NS}>
<office:automatic-styles>
{auto}
</office:automatic-styles>
<office:body>
<office:presentation>
{slides_xml}
<presentation:settings presentation:mouse-visible="false"/>
</office:presentation>
</office:body>
</office:document-content>"""

STYLES_XML = f"""<?xml version="1.0" encoding="UTF-8"?>
<office:document-styles {NS}>
<office:styles>
  <style:default-style style:family="graphic">
    <style:graphic-properties draw:fill="none" draw:stroke="none"/>
    <style:text-properties fo:font-family="Microsoft JhengHei" fo:font-size="18pt"/>
  </style:default-style>
  <style:default-style style:family="paragraph">
    <style:text-properties fo:font-family="Microsoft JhengHei" fo:font-size="18pt"/>
  </style:default-style>
</office:styles>
<office:master-styles>
  <draw:layer-set>
    <draw:layer draw:name="layout"/>
    <draw:layer draw:name="background"/>
    <draw:layer draw:name="backgroundobjects"/>
    <draw:layer draw:name="controls"/>
    <draw:layer draw:name="measurelines"/>
  </draw:layer-set>
  <style:master-page style:name="Default" style:display-name="Default"
    draw:style-name="dp1"/>
</office:master-styles>
</office:document-styles>"""

META_XML = """<?xml version="1.0" encoding="UTF-8"?>
<office:document-meta xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"
  xmlns:meta="urn:oasis:names:tc:opendocument:xmlns:meta:1.0" office:version="1.2">
<office:meta><meta:generator>pptx_to_odp.py</meta:generator></office:meta>
</office:document-meta>"""

MANIFEST_XML = """<?xml version="1.0" encoding="UTF-8"?>
<manifest:manifest xmlns:manifest="urn:oasis:names:tc:opendocument:xmlns:manifest:1.0"
  manifest:version="1.2">
<manifest:file-entry manifest:media-type="application/vnd.oasis.opendocument.presentation" manifest:full-path="/"/>
<manifest:file-entry manifest:media-type="text/xml" manifest:full-path="content.xml"/>
<manifest:file-entry manifest:media-type="text/xml" manifest:full-path="styles.xml"/>
<manifest:file-entry manifest:media-type="text/xml" manifest:full-path="meta.xml"/>
<manifest:file-entry manifest:media-type="text/xml" manifest:full-path="settings.xml"/>
</manifest:manifest>"""

SETTINGS_XML = """<?xml version="1.0" encoding="UTF-8"?>
<office:document-settings xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"
  office:version="1.2">
<office:settings/>
</office:document-settings>"""

# ── 主程式 ─────────────────────────────────────────────────────────
print(f"處理 {len(prs.slides)} 張投影片，尺寸 {SW:.1f}x{SH:.1f}cm")

all_slides = []
for i, slide in enumerate(prs.slides):
    print(f"  slide {i+1:02d}/{len(prs.slides)}", end='', flush=True)
    all_slides.append(slide_xml(i, slide))
    print(" OK")

slides_xml = '\n'.join(all_slides)
content    = build_content(slides_xml)

# 驗證 XML
from lxml import etree
try:
    etree.fromstring(content.encode('utf-8'))
    print("content.xml XML 驗證通過")
except etree.XMLSyntaxError as e:
    print(f"XML 錯誤: {e}")
    raise

with zipfile.ZipFile(OUTPUT, 'w', zipfile.ZIP_DEFLATED) as z:
    z.writestr('mimetype', 'application/vnd.oasis.opendocument.presentation',
               compress_type=zipfile.ZIP_STORED)
    z.writestr('META-INF/manifest.xml', MANIFEST_XML.encode('utf-8'))
    z.writestr('content.xml',  content.encode('utf-8'))
    z.writestr('styles.xml',   STYLES_XML.encode('utf-8'))
    z.writestr('meta.xml',     META_XML.encode('utf-8'))
    z.writestr('settings.xml', SETTINGS_XML.encode('utf-8'))

print(f"\nDone: {OUTPUT}")
print(f"段落樣式: {len(para_styles)}  文字樣式: {len(text_styles)}")
